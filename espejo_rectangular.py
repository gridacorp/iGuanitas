# espejo_rectangular.py

import os
from collections import Counter
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def read_text_file(path: str) -> list[str]:
    """Lee .txt o .py línea a línea."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return [l.rstrip("\n") for l in f]

def read_docx_file(path: str) -> list[str]:
    doc = Document(path)
    return [p.text for p in doc.paragraphs]

def read_pdf_file(path: str) -> list[str]:
    reader = PdfReader(path)
    return [page.extract_text() or "" for page in reader.pages]

def read_pptx_file(path: str) -> list[str]:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return texts

def compare_word_frequency(list1: list[str],
                           list2: list[str]) -> float:
    """% de diferencia en frecuencia de palabras."""
    words1 = " ".join(list1).lower().split()
    words2 = " ".join(list2).lower().split()

    c1 = Counter(words1)
    c2 = Counter(words2)

    all_words = set(c1) | set(c2)
    total_diff = sum(abs(c1[w] - c2[w]) for w in all_words)
    total_words = sum(c1.values()) + sum(c2.values())

    if total_words == 0:
        return 0.0
    return round((total_diff / total_words) * 100, 2)
