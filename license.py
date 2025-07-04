# espejo.py

import os
import pandas as pd
import numpy as np
from tkinter import filedialog, messagebox
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl

cancelar = False  # bandera global para cancelar proceso

def read_file(path):
    """Lee el contenido del archivo y devuelve DataFrame o lista."""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".xlsx", ".xls", ".xlsb", ".csv"):
        try:
            if ext == ".csv":
                return pd.read_csv(path, header=None, dtype=str, encoding="utf-8", error_bad_lines=False)
            else:
                return pd.read_excel(path, header=None, dtype=str, engine="openpyxl")
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")
    elif ext in (".txt", ".py"):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return [line.rstrip("\n") for line in f]
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")
    elif ext == ".docx":
        try:
            doc = Document(path)
            return [p.text for p in doc.paragraphs]
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")
    elif ext == ".pdf":
        try:
            r = PdfReader(path)
            return [page.extract_text() or "" for page in r.pages]
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")
    elif ext == ".pptx":
        try:
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return texts
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")
    else:
        raise Exception(f"Extensión no soportada: {ext}")

def compare_lists(list1, list2):
    """Compara dos listas línea a línea y devuelve porcentaje de diferencia."""
    max_len = max(len(list1), len(list2))
    if max_len == 0:
        return 0.0
    diffs = sum(1 for a, b in zip(list1, list2) if a != b)
    diffs += abs(len(list1) - len(list2))
    return round((diffs / max_len) * 100, 2)

def compare_files(files, status_label, progress_bar):
    """Compara todos los pares de archivos, muestra resumen y exporta CSV."""
    global cancelar
    cancelar = False

    if len(files) < 2:
        messagebox.showerror("Error", "Selecciona al menos dos archivos para comparar.")
        return

    total = (len(files) * (len(files) - 1)) // 2
    current = 0
    results = []

    for i in range(len(files)):
        for j in range(i + 1, len(files)):
            if cancelar:
                status_label.config(text="Proceso cancelado.")
                return

            f1, f2 = files[i], files[j]
            name1, name2 = os.path.basename(f1), os.path.basename(f2)
            try:
                content1 = read_file(f1)
                content2 = read_file(f2)

                if isinstance(content1, pd.DataFrame) and isinstance(content2, pd.DataFrame):
                    # Reindex para igualar forma
                    mr = max(content1.shape[0], content2.shape[0])
                    mc = max(content1.shape[1], content2.shape[1])
                    df1 = content1.reindex(index=range(mr), columns=range(mc), fill_value=np.nan)
                    df2 = content2.reindex(index=range(mr), columns=range(mc), fill_value=np.nan)
                    diff_mask = (df1 != df2) & ~(df1.isna() & df2.isna())
                    pct_diff = round((diff_mask.sum().sum() / (mr * mc)) * 100, 2)
                else:
                    pct_diff = compare_lists(content1, content2)

                results.append({
                    "Archivo 1": name1,
                    "Archivo 2": name2,
                    "% Diferencia": pct_diff
                })
            except Exception as e:
                results.append({
                    "Archivo 1": name1,
                    "Archivo 2": name2,
                    "% Diferencia": f"Error: {e}"
                })

            current += 1
            progress = (current / total) * 100
            progress_bar['value'] = progress
            status_label.config(text=f"Comparndo iguanas {current}/{total}...")
            status_label.update_idletasks()
            progress_bar.update_idletasks()

    # Mostrar resumen
    summary = "\n".join(
        f"{r['Archivo 1']} vs {r['Archivo 2']} → {r['% Diferencia']}%"
        for r in results
    )
    messagebox.showinfo("Resultado de comparación", summary)

    # Guardar CSV
    save_path = filedialog.asksaveasfilename(
        defaultextension=".csv",
        filetypes=[("CSV","*.csv")],
        title="Guardar resultados como CSV"
    )
    if save_path:
        df = pd.DataFrame(results)
        df.to_csv(save_path, index=False, encoding="utf-8")
        messagebox.showinfo("Guardado", f"Resultados guardados en:\n{save_path}")

    status_label.config(text="Comparación completada.")
    progress_bar['value'] = 100
