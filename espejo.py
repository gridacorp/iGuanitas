#espejo.py
import os 
import pandas as pd
import numpy as np
from tkinter import filedialog, messagebox
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from collections import Counter
import openpyxl

cancelar = False  # bandera global para cancelar proceso

def read_file(path):
    """Lee el contenido del archivo y devuelve:
    - dict de DataFrames para Excel/CSV (hojas)
    - lista de líneas para texto, pdf, docx, pptx"""
    ext = os.path.splitext(path)[1].lower()
    base = os.path.splitext(os.path.basename(path))[0]

    if ext == ".csv":
        try:
            try:
                df = pd.read_csv(
                    path,
                    header=None,
                    dtype=str,
                    encoding="utf-8",
                    sep=None,
                    engine="python",
                    on_bad_lines="skip"
                )
            except UnicodeDecodeError:
                df = pd.read_csv(
                    path,
                    header=None,
                    dtype=str,
                    encoding="latin-1",
                    sep=None,
                    engine="python",
                    on_bad_lines="skip"
                )
            return {base: df}
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")

    elif ext in (".xlsx", ".xls", ".xlsb"):
        try:
            # Quitar engine explícito para dejar que pandas elija el más rápido
            sheets = pd.read_excel(
                path,
                sheet_name=None,
                header=None,
                dtype=str
            )
            return sheets
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")

    elif ext in (".txt", ".py"):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return [line.rstrip("\n") for line in f]
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")

    elif ext == ".docx":
        try:
            doc = Document(path)
            return [p.text for p in doc.paragraphs]
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")

    elif ext == ".pdf":
        try:
            reader = PdfReader(path)
            return [page.extract_text() or "" for page in reader.pages]
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")

    elif ext == ".pptx":
        try:
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return texts
        except Exception as e:
            raise Exception(f"Error leyendo {os.path.basename(path)}: {e}")

    else:
        raise Exception(f"Extensión no soportada: {ext}")


def compare_lists(list1, list2):
    """Compara dos listas línea a línea y devuelve porcentaje de diferencia."""
    max_len = max(len(list1), len(list2))
    if max_len == 0:
        return 0.0
    diffs = sum(1 for a, b in zip(list1, list2) if a != b)
    diffs += abs(len(list1) - len(list2))
    return round((diffs / max_len) * 100, 2)


def compare_char_frequency(text1, text2):
    """Compara la frecuencia de caracteres entre dos listas de líneas y devuelve porcentaje de diferencia."""
    c1 = Counter(''.join(text1).lower())
    c2 = Counter(''.join(text2).lower())
    all_chars = set(c1.keys()).union(c2.keys())
    total_diff = sum(abs(c1.get(ch, 0) - c2.get(ch, 0)) for ch in all_chars)
    total_chars = sum(c1.values()) + sum(c2.values())
    if total_chars == 0:
        return 0.0
    return round((total_diff / total_chars) * 100, 2)


def compare_files(files, status_label, progress_bar):
    """Compara archivos:
    - Excel/CSV: % Diferencia por celdas (célula a célula)
    - Texto/Word/PDF/PPTX: % diferencia por líneas y caracteres."""
    global cancelar
    cancelar = False

    if len(files) < 2:
        messagebox.showerror("Error", "Selecciona al menos dos archivos para comparar.")
        return

    # Leer todos los archivos
    all_content = {}
    for f in files:
        try:
            all_content[f] = read_file(f)
        except Exception as e:
            all_content[f] = f"Error: {e}"

    # Calcular total de pares para la barra de progreso
    items = list(all_content.items())
    total_pairs = 0
    for i in range(len(items)):
        for j in range(i+1, len(items)):
            c1, c2 = items[i][1], items[j][1]
            if isinstance(c1, dict) and isinstance(c2, dict):
                total_pairs += len(c1) * len(c2)
            else:
                total_pairs += 1
    if total_pairs == 0:
        messagebox.showerror("Error", "No hay elementos válidos para comparar.")
        return

    results = []
    current = 0

    # Comparación par a par
    for i in range(len(items)):
        f1, c1 = items[i]
        for j in range(i+1, len(items)):
            f2, c2 = items[j]

            # Excel/CSV: comparación célula a célula
            if isinstance(c1, dict) and isinstance(c2, dict):
                for name1, df1 in c1.items():
                    for name2, df2 in c2.items():
                        # Igualar tamaño y rellenar con ""
                        mr = max(df1.shape[0], df2.shape[0])
                        mc = max(df1.shape[1], df2.shape[1])
                        d1 = df1.reindex(index=range(mr), columns=range(mc), fill_value="").astype(str)
                        d2 = df2.reindex(index=range(mr), columns=range(mc), fill_value="").astype(str)

                        # DEBUG opcional
                        print(f"[DEBUG] Comparando {name1} ({mr}×{mc}) vs {name2} ({mr}×{mc})")

                        # Vectorizar diff sobre numpy arrays
                        diff = (d1.values != d2.values)
                        pct = round(diff.sum() / (mr * mc) * 100, 2)

                        results.append({
                            "Archivo 1": os.path.basename(f1),
                            "Hoja 1": name1,
                            "Archivo 2": os.path.basename(f2),
                            "Hoja 2": name2,
                            "% Diferencia por celdas": pct,
                            "% Diferencia por palabras": "N/A"
                        })

                        current += 1
                        progress = current / total_pairs * 100
                        progress_bar['value'] = progress
                        status_label.config(text=f"Comparando {current}/{total_pairs}...")
                        status_label.update()
                        progress_bar.update()

            else:
                # Texto/Word/PDF/PPTX: comparación por líneas y caracteres
                cont1 = c1 if isinstance(c1, list) else []
                cont2 = c2 if isinstance(c2, list) else []
                pct_line = compare_lists(cont1, cont2)
                pct_char = compare_char_frequency(cont1, cont2)

                results.append({
                    "Archivo 1": os.path.basename(f1),
                    "Hoja 1": "-",
                    "Archivo 2": os.path.basename(f2),
                    "Hoja 2": "-",
                    "% Diferencia por líneas": pct_line,
                    "% Diferencia por palabras": pct_char
                })

                current += 1
                progress = current / total_pairs * 100
                progress_bar['value'] = progress
                status_label.config(text=f"Comparando {current}/{total_pairs}...")
                status_label.update()
                progress_bar.update()

    # Mostrar resultados y opción de guardar
    summary = "\n".join(
        f"{r['Archivo 1']} ({r.get('Hoja 1','-')}) vs {r['Archivo 2']} ({r.get('Hoja 2','-')}) → "
        + (f"Celdas: {r['% Diferencia por celdas']}%" if r.get('% Diferencia por celdas') != "N/A"
           else f"Líneas: {r['% Diferencia por líneas']}% | Letras: {r['% Diferencia por palabras']}%")
        for r in results
    )
    messagebox.showinfo("Resultado de comparación", summary)

    save_path = filedialog.asksaveasfilename(
        defaultextension=".csv",
        filetypes=[("CSV","*.csv")],
        title="Guardar resultados como CSV"
    )
    if save_path:
        pd.DataFrame(results).to_csv(save_path, index=False, encoding="utf-8")
        messagebox.showinfo("Guardado", f"Resultados guardados en:\n{save_path}")

    status_label.config(text="Comparación completada.")
    progress_bar['value'] = 100
