# comparador.py

import os                      # Operaciones con el sistema de archivos
import datetime                # Manejo de fechas y horas
import hashlib                 # Cálculo de hashes (SHA-256)
import threading               # Hilos para no bloquear la interfaz
import tkinter as tk           # Interfaz gráfica básica
from tkinter import filedialog, messagebox, ttk  # Diálogos, mensajes y widgets avanzados

import pandas as pd            # Lectura y manejo de datos en tablas
import numpy as np             # Operaciones numéricas
from docx import Document      # Lectura de archivos .docx
from PyPDF2 import PdfReader   # Lectura de archivos PDF
from pptx import Presentation  # Lectura de archivos PowerPoint
import openpyxl                # Lectura de archivos Excel .xlsx
import olefile                 # Lectura de metadatos de archivos .xls
import docx                    # Metadatos de archivos .docx
from pdf2image import convert_from_path  # Convertir PDF a imágenes
import pytesseract             # OCR para extraer texto de imágenes

from license import LICENSE_TEXT  # Importa la licencia desde license.py

# ----------------------
# Bandera global de cancelación
# ----------------------
cancelar = False

# ----------------------
# Función: Mostrar la licencia y esperar aceptación
# ----------------------
def show_license_and_get_acceptance():
    """Muestra un diálogo con el texto de la licencia y devuelve True si el usuario la acepta."""
    dlg = tk.Tk()
    dlg.title("Licencia de Uso")
    dlg.geometry("700x500")

    # Widget de texto para mostrar la licencia
    txt = tk.Text(dlg, wrap="word")
    txt.insert("1.0", LICENSE_TEXT)
    txt.config(state="disabled")
    txt.pack(fill=tk.BOTH, expand=True)

    accepted = tk.BooleanVar(value=False)

    def on_accept():
        accepted.set(True)
        dlg.destroy()

    def on_decline():
        dlg.destroy()

    # Botones Acepto / No acepto
    btn_frame = tk.Frame(dlg)
    tk.Button(btn_frame, text="Acepto",   command=on_accept,   width=12).pack(side=tk.LEFT, padx=5, pady=5)
    tk.Button(btn_frame, text="No acepto", command=on_decline,  width=12).pack(side=tk.LEFT, padx=5, pady=5)
    btn_frame.pack()

    dlg.mainloop()
    return accepted.get()

# ----------------------
# Función: Extraer texto de PDFs basados en imágenes
# ----------------------
def extract_text_from_image_pdf(file_path):
    """Convierte cada página del PDF a imagen y aplica OCR."""
    images = convert_from_path(file_path)
    return [pytesseract.image_to_string(img) for img in images]

# ----------------------
# Función: Lectura genérica de archivo según extensión
# ----------------------
def read_file(file_path):
    """
    Lee y extrae el contenido de un archivo.
    Retorna lista de líneas (texto) o DataFrame (tablas).
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.xlsx':
            return pd.read_excel(file_path, dtype=str, header=None, engine="openpyxl")
        if ext == '.xls':
            return pd.read_excel(file_path, dtype=str, header=None, engine="xlrd")
        if ext == '.xlsb':
            return pd.read_excel(file_path, dtype=str, header=None, engine="pyxlsb")
        if ext == '.csv':
            try:
                return pd.read_csv(file_path, dtype=str, header=None, encoding="utf-8")
            except UnicodeDecodeError:
                return pd.read_csv(file_path, dtype=str, header=None, encoding="latin1")
        if ext in ('.txt', '.py'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return [line.rstrip('\n') for line in f]
        if ext in ('.doc', '.docx'):
            doc = Document(file_path)
            return [p.text for p in doc.paragraphs]
        if ext == '.pdf':
            reader = PdfReader(file_path)
            pages = [p.extract_text() or "" for p in reader.pages]
            # Si no extrae texto, usa OCR
            return pages if any(pages) else extract_text_from_image_pdf(file_path)
        if ext in ('.ppt', '.pptx'):
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts += shape.text.split('\n')
            return texts
        raise ValueError(f"Extensión no soportada: {ext}")
    except Exception as e:
        raise RuntimeError(f"No se pudo leer {file_path}: {e}")

# ----------------------
# Función: Comparar dos listas de texto
# ----------------------
def compare_lists(a, b):
    """
    Compara línea a línea dos listas de texto,
    devuelve porcentaje de líneas diferentes.
    """
    max_len = max(len(a), len(b))
    a += [""] * (max_len - len(a))
    b += [""] * (max_len - len(b))
    diffs = sum(1 for x, y in zip(a, b) if x != y)
    return round((diffs / max_len) * 100, 2) if max_len else 0.0

# ----------------------
# Función: Comparar múltiples archivos y mostrar progreso
# ----------------------
def compare_files(files, status_label, progress_bar):
    """
    Compara cada par de archivos en 'files'.
    Actualiza etiqueta y barra de progreso.
    Guarda resultados en CSV al terminar.
    """
    global cancelar
    cancelar = False
    if len(files) < 2:
        messagebox.showerror("Error", "Selecciona al menos dos archivos para comparar.")
        return

    total = (len(files) * (len(files) - 1)) // 2
    current = 0
    results = []

    for i in range(len(files)):
        if cancelar:
            status_label.config(text="Proceso cancelado.")
            return
        a = files[i]
        try:
            data_a = read_file(a)
        except Exception as e:
            messagebox.showerror("Error", str(e))
            continue

        for j in range(i + 1, len(files)):
            if cancelar:
                status_label.config(text="Proceso cancelado.")
                return
            b = files[j]
            try:
                data_b = read_file(b)
            except Exception as e:
                messagebox.showerror("Error", str(e))
                continue

            current += 1
            # Actualiza estado y progreso
            status_label.config(text=f"Revisando {current}/{total}\n'{os.path.basename(a)}' vs '{os.path.basename(b)}'")
            progress_bar["value"] = (current / total) * 100
            status_label.update_idletasks()
            progress_bar.update_idletasks()

            # Si es DataFrame (hoja de cálculo), compara celda a celda
            if isinstance(data_a, pd.DataFrame):
                df1, df2 = data_a.copy(), data_b.copy()
                mr, mc = max(df1.shape[0], df2.shape[0]), max(df1.shape[1], df2.shape[1])
                df1 = df1.reindex(index=range(mr), columns=range(mc), fill_value=np.nan)
                df2 = df2.reindex(index=range(mr), columns=range(mc), fill_value=np.nan)
                ch = (df1 != df2) & ~(df1.isna() & df2.isna())
                pct = round((ch.sum().sum() / (mr * mc)) * 100, 2)
            else:
                pct = compare_lists(data_a, data_b)

            results.append([os.path.basename(a), os.path.basename(b), pct])

    # Guardar resultados como CSV
    save = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV", "*.csv")])
    if save:
        pd.DataFrame(results, columns=["Archivo 1", "Archivo 2", "% Diferencia"]) \
          .to_csv(save, index=False, encoding="utf-8")
        messagebox.showinfo("Listo", f"Guardado en {save}")

    status_label.config(text="Comparación completada.")
    progress_bar["value"] = 100

# ----------------------
# Función: Calcular hash SHA-256 de un archivo
# ----------------------
def calcular_hash(path):
    """Lee el archivo en bloques y devuelve su SHA-256."""
    try:
        h = hashlib.sha256()
        with open(path, 'rb') as f:
            while chunk := f.read(8192):
                h.update(chunk)
        return h.hexdigest()
    except:
        return "Error"

# ----------------------
# Función: Extraer metadatos de distintos tipos de archivo
# ----------------------
def extraer_metadatos(archivo):
    """
    Retorna lista [autor, último editor, fecha creación, fecha modificación]
    según el tipo de archivo.
    """
    ext = os.path.splitext(archivo)[1].lower()
    try:
        if ext == '.xlsx':
            wb = openpyxl.load_workbook(archivo, data_only=True)
            p = wb.properties
            return [
                p.creator or "-", p.lastModifiedBy or "-",
                p.created.strftime("%Y-%m-%d %H:%M:%S") if p.created else "-",
                p.modified.strftime("%Y-%m-%d %H:%M:%S") if p.modified else "-"
            ]
        if ext == '.xls':
            ole = olefile.OleFileIO(archivo)
            m = ole.get_metadata()
            return [
                m.author or "-", m.last_saved_by or "-",
                m.create_time.strftime("%Y-%m-%d %H:%M:%S") if m.create_time else "-",
                m.last_saved_time.strftime("%Y-%m-%d %H:%M:%S") if m.last_saved_time else "-"
            ]
        if ext == '.docx':
            d = docx.Document(archivo)
            p = d.core_properties
            return [
                p.author or "-", p.last_modified_by or "-",
                p.created.strftime("%Y-%m-%d %H:%M:%S") if p.created else "-",
                p.modified.strftime("%Y-%m-%d %H:%M:%S") if p.modified else "-"
            ]
        if ext == '.doc':
            return ["No soportado"] * 4
        if ext == '.pdf':
            r = PdfReader(archivo)
            i = r.metadata
            return [i.author or "-", i.producer or "-", i.creation_date or "-", i.modification_date or "-"]
        if ext == '.pptx':
            pr = Presentation(archivo)
            p = pr.core_properties
            return [
                p.author or "-", p.last_modified_by or "-",
                p.created.strftime("%Y-%m-%d %H:%M:%S") if p.created else "-",
                p.modified.strftime("%Y-%m-%d %H:%M:%S") if p.modified else "-"
            ]
        if ext == '.ppt':
            return ["No soportado"] * 4
        # Texto plano y CSV: usa fechas del sistema de archivos
        stat = os.stat(archivo)
        c = datetime.datetime.fromtimestamp(stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S")
        m = datetime.datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
        return ["N/A", "N/A", c, m]
    except:
        return ["Error"] * 4

# ----------------------
# Función: Mostrar metadatos en una tabla de la GUI
# ----------------------
def mostrar_metadatos(root):
    files = filedialog.askopenfilenames(
        filetypes=[("Soportados", "*.xlsx *.xls *.xlsb *.csv *.txt *.py *.docx *.doc *.pdf *.pptx *.ppt")]
    )
    if not files:
        return

    win = tk.Toplevel(root)
    win.title("Metadatos y Hash")
    cols = ["Archivo", "Creador", "Editor", "Creación", "Modificación", "SHA-256"]
    tv = ttk.Treeview(win, columns=cols, show='headings')
    for c in cols:
        tv.heading(c, text=c)
        tv.column(c, width=130, anchor="center")
    tv.pack(fill=tk.BOTH, expand=True)

    data = []
    for f in files:
        m = extraer_metadatos(f)
        h = calcular_hash(f)
        row = [os.path.basename(f), *m, h]
        tv.insert("", tk.END, values=row)
        data.append(row)

    # Opción de guardar CSV con metadatos
    if messagebox.askyesno("Guardar CSV", "¿Guardar metadatos en CSV?"):
        path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV", "*.csv")])
        if path:
            pd.DataFrame(data, columns=cols).to_csv(path, index=False, encoding="utf-8")
            messagebox.showinfo("Guardado", f"CSV guardado en {path}")

# ----------------------
# Función: Seleccionar archivos y lanzar hilo de comparación
# ----------------------
def pick_files(tipo, exts, status_label, progress_bar):
    files = filedialog.askopenfilenames(filetypes=[(tipo, exts)])
    if files:
        threading.Thread(
            target=compare_files,
            args=(files, status_label, progress_bar),
            daemon=True
        ).start()

# ----------------------
# Función: Marcar cancelación del proceso
# ----------------------
def cancelar_proceso():
    global cancelar
    cancelar = True

# ----------------------
# Función principal: arma la interfaz y coordina todo
# ----------------------
def main():
    # Mostrar licencia y validar aceptación
    if not show_license_and_get_acceptance():
        return

    root = tk.Tk()
    root.title("Comparador de Archivos")
    root.geometry("600x620")

    tk.Label(
        root,
        text="Selecciona qué tipo de documento comparar:",
        font=("Arial", 12)
    ).pack(pady=10)

    status_label = tk.Label(
        root, text="", fg="blue", font=("Arial", 10),
        wraplength=580, justify="center"
    )
    status_label.pack(pady=5)

    progress_bar = ttk.Progressbar(root, length=550, mode='determinate')
    progress_bar.pack(pady=5)

    # Botones para cada tipo de archivo
    btn_specs = [
        ("📊 Hojas de cálculo", "*.xlsx *.xls *.xlsb *.csv"),
        ("📄 Texto plano",      "*.txt *.py"),
        ("📝 Word",             "*.doc *.docx"),
        ("📕 PDF",              "*.pdf"),
        ("📽 PowerPoint",       "*.ppt *.pptx"),
    ]
    for text, exts in btn_specs:
        tk.Button(
            root,
            text=f"{text} ({exts})",
            width=60,
            command=lambda t=text, e=exts: pick_files(t, e, status_label, progress_bar)
        ).pack(pady=3)

    # Botones adicionales: propiedades/hash y cancelar
    tk.Button(
        root,
        text="🔍 Comparar Propiedades y Hash",
        width=60,
        command=lambda: mostrar_metadatos(root)
    ).pack(pady=15)

    tk.Button(
        root,
        text="❌ Cancelar proceso",
        bg="red", fg="white",
        width=60,
        command=cancelar_proceso
    ).pack(pady=5)

    root.mainloop()

if __name__ == "__main__":
    main()

